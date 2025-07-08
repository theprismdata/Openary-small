import json
import os
import time
from pprint import pprint
import pymongo
import pymysql
import yaml
import pika
import pdfplumber
import logging
from logging.handlers import TimedRotatingFileHandler
from langchain.document_loaders import UnstructuredWordDocumentLoader
from langchain.document_loaders import TextLoader
import pandas as pd
import re
from HwpParser import HWPExtractor
from pptx import Presentation
from minio import Minio
from datetime import datetime
from pytz import timezone
from functools import lru_cache
from contextlib import contextmanager

print("build:2025-02-26#10:15")

# 환경 설정 및 설정 파일 로드
ENV = os.getenv('ENVIRONMENT', 'development')
config_file = f'../config/svc-set.{"debug." if ENV == "development" else ""}yaml'

# 설정 파일은 한 번만 로드
with open(config_file) as f:
    config = yaml.load(f, Loader=yaml.FullLoader)

# 로깅 설정
logger = logging.getLogger("Rotating Log")
logger.setLevel(logging.DEBUG)

f_format = logging.Formatter('[%(asctime)s][%(levelname)s|%(filename)s:%(lineno)s] --- %(message)s')

if not os.path.exists("log"):
    os.mkdir("log")
path = "log/preprocess.log"
file_handler = TimedRotatingFileHandler(path,
                                        when="h",
                                        interval=1,
                                        backupCount=24)
file_handler.namer = lambda name: name + ".txt"
file_handler.setFormatter(f_format)

stream_handler = logging.StreamHandler()
stream_handler.setFormatter(f_format)

logger.addHandler(file_handler)
logger.addHandler(stream_handler)

# 설정 값 추출
mqtt_info = config['mqtt']
mqtt_id = mqtt_info['id']
mqtt_pwd = mqtt_info['pwd']
mqtt_address = mqtt_info['address']
mqtt_port = mqtt_info['port']
mqtt_virtualhost = mqtt_info['virtualhost']

RABBITMQ_SVC_QUEUE = config['RABBITMQ_SVC_QUEUE']

mongo_host = config['database']['mongodb']['mongo_host']
mongo_port = config['database']['mongodb']['mongo_port']
mongo_user = config['database']['mongodb']['mongo_user']
mongo_passwd = config['database']['mongodb']['mongo_passwd']
auth_source = config['database']['mongodb']['auth_source']

minio_info = config['minio']
minio_address = minio_info['address']
accesskey = minio_info['accesskey']
secretkey = minio_info['secretkey']

opds_system_db = config['database']['opds_system_db']

mongo_uri = f"mongodb://{mongo_user}:{mongo_passwd}@{mongo_host}:{mongo_port}/?authSource={auth_source}&authMechanism=SCRAM-SHA-1"

# 큐 이름 설정
OPDS_SUMMARY_REQ_Qname = RABBITMQ_SVC_QUEUE['SUMMARY_Q_NAME']
OPDS_SUMMARY_REQ_ROUTE = RABBITMQ_SVC_QUEUE['SUMMARY_ROUTE_KEY']
OPDS_EMBEDDING_REQ_Qname = RABBITMQ_SVC_QUEUE['EMBEDDING_Q_NAME']
OPDS_EMBEDDING_ROUTE = RABBITMQ_SVC_QUEUE['EMBEDDING_ROUTE_KEY']


# 데이터베이스 연결 관리를 위한 컨텍스트 매니저들
@contextmanager
def get_mongo_connection():
    """MongoDB 연결 컨텍스트 매니저"""
    client = pymongo.MongoClient(mongo_uri, maxPoolSize=10)
    try:
        yield client[auth_source]
    finally:
        client.close()


@contextmanager
def get_mysql_connection():
    """MySQL 연결 컨텍스트 매니저"""
    conn = pymysql.connect(
        user=opds_system_db["id"],
        password=opds_system_db["pwd"],
        database=opds_system_db["database"],
        host=opds_system_db["address"],
        port=opds_system_db["port"]
    )
    try:
        yield conn
    finally:
        conn.close()


@contextmanager
def get_minio_client():
    """MinIO 클라이언트 컨텍스트 매니저"""
    client = Minio(
        minio_address,
        access_key=accesskey,
        secret_key=secretkey,
        secure=False
    )
    try:
        yield client
    finally:
        pass  # MinIO 클라이언트는 명시적으로 닫을 필요가 없음


# RabbitMQ 연결 관리 클래스
class RabbitMQConnection:
    def __init__(self, connection_name, credentials, params):
        self.connection_name = connection_name
        self.connection = pika.BlockingConnection(params)
        self.channel = self.connection.channel()

    def close(self):
        if self.connection and not self.connection.is_closed:
            self.connection.close()

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        self.close()


# 사용자 정보 조회 함수에 캐싱 적용
@lru_cache(maxsize=100)
def get_userid_fromcode(user_code):
    """사용자 코드로 사용자 이름 조회 (캐싱 적용)"""
    with get_mysql_connection() as conn:
        with conn.cursor() as cs:
            sql = f'SELECT `id`, name FROM tb_user WHERE user_code="{user_code}"'
            cs.execute(sql)
            rs = cs.fetchall()
            user_name_df = pd.DataFrame(rs, columns=['id', 'name'])

    if user_name_df.shape[0] == 1:
        user_name_df = user_name_df.iloc[0].to_dict()
        return user_name_df["name"]
    return None


@lru_cache(maxsize=100)
def get_usercode_from_email(user_email):
    """이메일로 사용자 코드 조회 (캐싱 적용)"""
    with get_mysql_connection() as conn:
        with conn.cursor() as cs:
            sql = f'SELECT `id`, user_code FROM tb_user WHERE email="{user_email}"'
            cs.execute(sql)
            rs = cs.fetchall()
            user_name_df = pd.DataFrame(rs, columns=['id', 'user_code'])

    if user_name_df.shape[0] == 1:
        user_name_df = user_name_df.iloc[0].to_dict()
        return user_name_df["user_code"]
    return None


def update_extract_progress(user_code, doc_id, progress):
    """추출 진행률 업데이트"""
    with get_mysql_connection() as conn:
        with conn.cursor() as cs:
            sql = f"""UPDATE {opds_system_db["database"]}.tb_llm_doc 
                      SET extract_page_rate = {progress} 
                      WHERE userid='{user_code}' AND id = '{doc_id}'"""
            cs.execute(sql)
            conn.commit()


def update_doc_status(doc_id, status="extracted", summary="분석중 입니다..."):
    """문서 상태 업데이트"""
    with get_mysql_connection() as conn:
        with conn.cursor() as cs:
            start_time = datetime.now(timezone('Asia/Seoul')).strftime("%Y-%m-%d %H:%M:%S.%f")
            sql = f"""UPDATE {opds_system_db["database"]}.tb_llm_doc 
                      SET status = "{status}", 
                          summary = "{summary}", 
                          process_start = "{start_time}" 
                      WHERE id = {doc_id}"""
            cs.execute(sql)
            conn.commit()
            logger.debug(sql)


def get_pdfpage_info_by_plumber(user_code, doc_id, pdf_sourcepath) -> str:
    """PDF 파일에서 텍스트 추출"""
    print(f"Source path {pdf_sourcepath}")
    try:
        # 임시 디렉토리 생성
        if not os.path.exists("tmp/minio_file/"):
            os.makedirs("tmp/minio_file/")
        downloaded_file_path = "tmp/minio_file/" + pdf_sourcepath

        # MinIO에서 파일 다운로드
        with get_minio_client() as client:
            client.fget_object(user_code, pdf_sourcepath, downloaded_file_path)

        # PDF 텍스트 추출
        with pdfplumber.open(downloaded_file_path) as pdfplumb_obj:
            file_extract_contents = ""
            total_pages = len(pdfplumb_obj.pages)

            for page_num, plumb_page in enumerate(pdfplumb_obj.pages):
                page_plumb_contents = {}
                logger.debug(f"Page {page_num} start")
                table_list = []

                # 테이블 추출
                for table_info in pdfplumb_obj.pages[page_num].find_tables():
                    x0, y0, x1, y1 = table_info.bbox
                    table_list.append((x0, y0, x1, y1))
                    table = table_info.extract()
                    df = pd.DataFrame(table[1::], columns=table[0])
                    df.replace('\x00', '', inplace=True)
                    df.replace('Ÿ', '*', inplace=True)
                    page_plumb_contents[int(y0)] = df.to_markdown()

                # 텍스트 추출
                for content in pdfplumb_obj.pages[page_num].extract_text_lines():
                    x0 = content['x0']
                    y0 = content['top']
                    x1 = content['x1']
                    y1 = content['bottom']
                    if len(table_list) > 0:
                        if table_list[0][0] < x0 and table_list[0][1] < y0 and table_list[0][2] > x1 and table_list[0][
                            3] > y1:
                            pass
                        else:
                            page_plumb_contents[int(y0)] = content['text']
                    else:
                        page_plumb_contents[int(y0)] = content['text']

                # 내용 조합
                if len(page_plumb_contents) > 0:
                    pos_list = sorted(page_plumb_contents.keys())
                    for position in pos_list:
                        file_extract_contents += page_plumb_contents[position] + "\n"
                logger.debug(f"Page {page_num} finish go next")

                # 10% 단위로만 진행률 업데이트
                page_rate = int(((page_num + 1) / total_pages) * 100)
                if page_rate % 10 == 0 or page_num == 0 or page_num == total_pages - 1:
                    update_extract_progress(user_code, doc_id, page_rate)

            # 텍스트 정제
            file_extract_contents = re.sub(r"(?<![\.\?\!])\n", " ", file_extract_contents)
            file_extract_contents = re.sub(r"\(cid:[0-9]+\)", "", file_extract_contents)
            logger.debug(f"{pdf_sourcepath} extracted")

            return file_extract_contents
    except Exception as e:
        logger.error(f"PDF 처리 오류: {str(e)}")
        print(e)
        return ""


def extract_content_from_bucket(user_code, doc_id, file_name):
    """다양한 파일 형식에서 콘텐츠 추출"""
    try:
        formed_clear_contents = ''

        # 파일 확장자에 따른 처리
        if file_name.lower().endswith('.pdf'):
            formed_clear_contents = get_pdfpage_info_by_plumber(user_code, doc_id, file_name)
            logger.debug("PDF extract success")

        elif file_name.endswith('.hwp'):
            hwp_obj = HWPExtractor(file_name)
            formed_clear_contents = hwp_obj.get_text()
            logger.debug("HWP extract success")

        elif file_name.lower().endswith(('.docx', '.doc')):
            logger.debug(f'word loader for {file_name}')
            loader = UnstructuredWordDocumentLoader(file_name)
            docs = loader.load()
            for page in docs:
                formed_clear_contents += page.page_content
            logger.debug("Word extract success")

        elif file_name.endswith('.txt'):
            loader = TextLoader(file_name)
            docs = loader.load()
            for page in docs:
                formed_clear_contents += page.page_content
            logger.debug("TXT extract success")

        elif file_name.lower().endswith(('.xlsx', '.xls')):
            logger.debug(f'excel loader for {file_name}')
            df = pd.read_excel(file_name)
            formed_clear_contents = df.to_markdown()
            logger.debug("Excel extract success")

        elif file_name.endswith('.csv'):
            logger.debug(f'csv loader for {file_name}')
            df = pd.read_csv(file_name)
            formed_clear_contents = df.to_markdown()
            logger.debug("CSV extract success")

        elif file_name.lower().endswith(('.pptx', '.ppt')):
            try:
                logger.debug(f'ppt(x) loader for {file_name}')
                prs = Presentation(file_name)
                for idx, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                formed_clear_contents += run.text + '\r\n'
                logger.debug("PPT extract success")
            except Exception as e:
                logger.error(f"PPT 처리 오류: {str(e)}")
                return 0, None, str(e)

        else:
            logger.error(f"Error: invalid file type {file_name}")
            return 0, None, f"Unsupported file type: {file_name}"

    except Exception as e:
        logger.error(f"Error processing {file_name}: {e}")
        return 0, None, str(e)

    return 1, formed_clear_contents, "ok"


def store_document_content(user_code, doc_id, file_name, clean_contents):
    """문서 내용을 MongoDB에 저장"""
    with get_mongo_connection() as mongodb_genai:
        user_collection = mongodb_genai[user_code]

        # 기존 문서 확인
        file_doc = user_collection.find_one({"filename": file_name, "id": doc_id}, {'_id': False})

        if file_doc is None:
            # 새 문서 삽입
            file_doc = {
                'filename': file_name,
                'id': doc_id,
                'clean_doc': clean_contents,
                'cleansing': 'finish',
                'summary': 'ready',
                'embedding': 'ready'
            }
            user_collection.insert_one(file_doc)
            logger.debug(f"mongo insert id {doc_id}")
        else:
            # 기존 문서 업데이트
            user_collection.update_one(
                {'filename': file_name, "id": doc_id},
                {'$set': {
                    'clean_doc': clean_contents,
                    'cleansing': 'finish',
                    'summary': 'ready',
                    'embedding': 'ready'
                }}
            )
            logger.debug(f"mongo update id {doc_id}")


def publish_to_queues(user_email, user_code, doc_id, file_name):
    """다른 서비스 큐에 메시지 발행"""
    msg_json = {
        "user_email": user_email,
        "user_code": user_code,
        "doc_id": doc_id,
        "file_name": file_name
    }
    message = json.dumps(msg_json)

    # RabbitMQ 연결 생성
    credentials = pika.PlainCredentials(mqtt_id, mqtt_pwd)
    param = pika.ConnectionParameters(
        mqtt_address,
        mqtt_port,
        mqtt_virtualhost,
        credentials,
        heartbeat=30,
        blocked_connection_timeout=30
    )

    # 요약 큐에 발행
    with RabbitMQConnection("summary", credentials, param) as summary_conn:
        summary_conn.channel.queue_declare(queue=OPDS_SUMMARY_REQ_Qname)
        summary_conn.channel.basic_publish(
            exchange='',
            routing_key=OPDS_SUMMARY_REQ_ROUTE,
            body=message
        )
        logger.debug("PREPROCESS TO SUMMARY QUEUE")

    # 임베딩 큐에 발행
    with RabbitMQConnection("embedding", credentials, param) as embedding_conn:
        embedding_conn.channel.queue_declare(queue=OPDS_EMBEDDING_REQ_Qname)
        embedding_conn.channel.basic_publish(
            exchange='',
            routing_key=OPDS_EMBEDDING_ROUTE,
            body=message
        )
        logger.debug("PREPROCESS TO EMBEDDING QUEUE")


def wait_mq_signal(ch, method, properties, body):
    """메시지 큐 메시지 처리 콜백"""
    body = body.decode('utf-8')
    msg_json = json.loads(body)
    logger.info("PREPROCESS SUBSCRIBE")
    logger.debug(msg_json)

    try:
        # 사용자 정보 및 파일 정보 추출
        user_email = msg_json['user_email']
        user_code = get_usercode_from_email(user_email)
        doc_id = msg_json['doc_id']
        file_name = msg_json['file_name']

        if not user_code:
            logger.error(f"User code not found for email: {user_email}")
            ch.basic_ack(delivery_tag=method.delivery_tag)
            return

        # 파일 내용 추출
        sts, clean_contents, msg = extract_content_from_bucket(user_code, doc_id, file_name)

        if clean_contents is not None:
            # MongoDB에 문서 내용 저장
            store_document_content(user_code, doc_id, file_name, clean_contents)

            # 문서 상태 업데이트
            update_doc_status(doc_id)

            # 다른 서비스 큐에 메시지 발행
            publish_to_queues(user_email, user_code, doc_id, file_name)

            # 메시지 처리 완료 확인
            ch.basic_ack(delivery_tag=method.delivery_tag)
            logger.debug("Message processing completed and acknowledged")
        else:
            logger.error(f"Failed to extract content from {file_name}: {msg}")
            ch.basic_ack(delivery_tag=method.delivery_tag)

    except Exception as e:
        logger.error(f"Error processing message: {str(e)}")
        ch.basic_ack(delivery_tag=method.delivery_tag)  # 오류가 있어도 메시지 승인 (무한 재시도 방지)


def create_connection():
    """RabbitMQ 연결 생성"""
    credentials = pika.PlainCredentials(mqtt_id, mqtt_pwd)
    param = pika.ConnectionParameters(
        host=mqtt_address,
        port=mqtt_port,
        virtual_host=mqtt_virtualhost,
        credentials=credentials,
        heartbeat=30,
        blocked_connection_timeout=30
    )
    return pika.BlockingConnection(param)


def cleanup_connection(connection):
    """연결 정리 함수"""
    try:
        if connection and not connection.is_closed:
            connection.close()
    except pika.exceptions.ConnectionWrongStateError:
        logger.warning("Connection already closed")
    except Exception as e:
        logger.error(f"Error closing connection: {str(e)}")


if __name__ == '__main__':
    while True:
        connection = None
        try:
            logger.info("Starting RabbitMQ connection")
            connection = create_connection()

            OPDS_PREP_SUB_Channel = connection.channel()
            OPDS_RREP_REQ_Q = RABBITMQ_SVC_QUEUE['PREPROCES_Q_NAME']
            OPDS_PREP_SUB_Channel.queue_declare(queue=OPDS_RREP_REQ_Q)
            OPDS_PREP_SUB_Channel.basic_qos(prefetch_count=1)  # 한 번에 하나의 메시지만 처리

            OPDS_PREP_SUB_Channel.basic_consume(
                queue=OPDS_RREP_REQ_Q,
                on_message_callback=wait_mq_signal,
                auto_ack=False
            )

            print("Preprocess consumer start")
            logger.info("Preprocess consumer start")
            OPDS_PREP_SUB_Channel.start_consuming()

        except pika.exceptions.StreamLostError:
            logger.error("Connection lost. Reconnecting...")
            time.sleep(5)
            cleanup_connection(connection)

        except pika.exceptions.ConnectionClosedByBroker:
            logger.error("Connection closed by broker. Reconnecting...")
            time.sleep(5)
            cleanup_connection(connection)

        except KeyboardInterrupt:
            logger.info("Preprocess consumer stopped by user")
            cleanup_connection(connection)
            break

        except Exception as e:
            logger.error(f"Unexpected error: {str(e)}")
            time.sleep(5)
            cleanup_connection(connection)